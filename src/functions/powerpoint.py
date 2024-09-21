import os
from pptx import Presentation
from pptx.util import Inches
from config import CHARTS_FOLDER, REPORTS_FOLDER

def create_powerpoint(charts_folder=CHARTS_FOLDER, reports_folder=REPORTS_FOLDER, ppt_filename='financial_report.pptx'):
    prs = Presentation()

    for chart_file in os.listdir(charts_folder):
        if chart_file.endswith('.png') or chart_file.endswith('.jpg'):
            slide_layout = prs.slide_layouts[5]  # Blank slide
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            title_text = os.path.splitext(chart_file)[0].replace('_', ' ').title()
            title_shape.text = title_text

            img_path = os.path.join(charts_folder, chart_file)
            left = Inches(1)
            top = Inches(1.5)
            height = Inches(5.5)
            slide.shapes.add_picture(img_path, left, top, height=height)

    ppt_path = os.path.join(reports_folder, ppt_filename)
    prs.save(ppt_path)
    return ppt_path
